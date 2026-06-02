"""
生成介绍PPT v3 — 对齐参考PPT风格（Poppins + Noto Sans SC, 青蓝主题）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

prs = Presentation()
prs.slide_width = Emu(12192000)   # 16:9
prs.slide_height = Emu(6858000)

# ── 颜色 ──
BG       = RGBColor(0xF8, 0xFA, 0xFC)
DARK     = RGBColor(0x0F, 0x17, 0x2A)
WHITE    = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT   = RGBColor(0x06, 0xB6, 0xD4)
BODY     = RGBColor(0x33, 0x41, 0x55)
MUTED    = RGBColor(0x64, 0x74, 0x8B)
SUBTLE   = RGBColor(0x94, 0xA3, 0xB8)
RED      = RGBColor(0xEF, 0x44, 0x44)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
CARD_BG  = RGBColor(0xF1, 0xF5, 0xF9)
DARK_CARD= RGBColor(0x1E, 0x29, 0x3B)
LINE     = RGBColor(0xE2, 0xE8, 0xF0)

FONT_EN = 'Poppins'
FONT_CN = 'Noto Sans SC'

def add_bg(slide, color=BG):
    """设置纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    """添加矩形"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_name=FONT_CN, size=Pt(12), color=BODY, bold=False, align=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multi_text(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT):
    """多行文本 (lines: [(text, font_name, size, color, bold), ...])"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_name, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = font_name
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def slide_title(slide, text, y=Cm(0.6)):
    add_text_box(slide, Cm(1.2), y, Cm(22), Cm(1.2), text, FONT_EN, Pt(28.5), DARK, True)
    add_shape(slide, Cm(1.2), y + Cm(1.1), Cm(3), Pt(3), ACCENT)

def card(slide, x, y, w, h, title_en, body_cn, accent=ACCENT):
    """内容卡片"""
    add_shape(slide, x, y, w, h, CARD_BG)
    # 顶部色条
    add_shape(slide, x, y, w, Pt(3), accent)
    add_text_box(slide, x + Cm(0.5), y + Cm(0.3), w - Cm(1), Cm(0.7), title_en, FONT_EN, Pt(15), DARK, True)
    add_text_box(slide, x + Cm(0.5), y + Cm(1.0), w - Cm(1), h - Cm(1.3), body_cn, FONT_CN, Pt(11), MUTED)


# ══════════════════════════════════════════════════
# SLIDE 1: 封面
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(s, DARK)

# 装饰线
add_shape(s, Cm(1.5), Cm(2.5), Cm(3), Pt(3), ACCENT)

add_multi_text(s, Cm(1.5), Cm(3.0), Cm(22), Cm(4.5), [
    ("基于AI规则引擎的\n内部存货交易未实现利润\n自动抵消系统", FONT_EN, Pt(39), WHITE, True),
])
add_text_box(s, Cm(1.5), Cm(9.8), Cm(22), Cm(0.8), 'DRP系统 T11 合并报表与一键出表 — 子模块', FONT_CN, Pt(14), SUBTLE)

add_multi_text(s, Cm(1.5), Cm(11.5), Cm(22), Cm(3), [
    ('中国社会科学院大学 · 2025级会计专硕二班', FONT_CN, Pt(13), WHITE, True),
    ('小组成员：尚浩然、李金阳 | 指导老师：张金昌', FONT_CN, Pt(12), SUBTLE, False),
    ('2026年6月', FONT_CN, Pt(11), MUTED, False),
])


# ══════════════════════════════════════════════════
# SLIDE 2: 背景与痛点
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, '项目背景与业务痛点')

# 左侧-政策
add_text_box(s, Cm(1.2), Cm(2.2), Cm(11), Cm(0.6), '政策与业务双重驱动', FONT_EN, Pt(18), DARK, True)

bullets_policy = [
    ('国资委1号文：应用AI大模型、OCR技术实现管理智能化。', ACCENT, True),
    ('国资委2号文：数据自动采集、模型自动分析、风险自动预警。', ACCENT, True),
    ('CAS33：合并报表必须抵消内部交易未实现利润。', ACCENT, True),
    ('集团型企业涉及数十家法人主体，月度内部交易数百笔。', BODY, False),
    ('交易类型复杂（顺流/逆流/平流），抵消逻辑各不相同。', BODY, False),
]
add_multi_text(s, Cm(1.2), Cm(3.0), Cm(11), Cm(4), 
    [("• " + t, FONT_CN, Pt(12), c, b) for t, c, b in bullets_policy])

# 右侧-四大痛点卡片
add_text_box(s, Cm(13), Cm(2.2), Cm(11), Cm(0.6), '传统手工抵消的四大痛点', FONT_EN, Pt(18), DARK, True)

pain_cards = [
    (Cm(13), Cm(3.0), '① 工作量巨大', '数十家主体、数百笔交易，\nExcel逐笔勾对耗时漫长', RED),
    (Cm(17.5), Cm(3.0), '② 极易出错', '顺流、逆流、平流抵消比例\n和计算逻辑各不相同', AMBER),
    (Cm(13), Cm(6.2), '③ 审计追溯困难', '年结审计时无法完整还原\n每笔抵消的底层计算过程', AMBER),
    (Cm(17.5), Cm(6.2), '④ 监管要求脱节', '纯手工无法满足国资委\n"穿透式监管"要求', RED),
]
for x, y, title, body, accent in pain_cards:
    card(s, x, y, Cm(4.2), Cm(2.8), title, body, accent)


# ══════════════════════════════════════════════════
# SLIDE 3: 解决方案
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, '解决方案：AI规则引擎与自动化')

# 核心流程
add_text_box(s, Cm(1.2), Cm(2.2), Cm(11), Cm(0.6), '核心处理流程', FONT_EN, Pt(18), DARK, True)
flow = ['① 上传Excel三表', '② 智能识别交易类型', '③ 自动计算未实现利润', '④ JSON规则引擎生成分录', '⑤ 试算平衡借贷验证', '⑥ 导出报告Excel/PDF']
for i, step in enumerate(flow):
    y = Cm(3.0) + Cm(0.9) * i
    add_shape(s, Cm(1.5), y, Cm(0.6), Cm(0.6), ACCENT)
    add_text_box(s, Cm(2.4), y - Cm(0.05), Cm(9), Cm(0.7), step, FONT_CN, Pt(12), BODY)

# 全场景覆盖
add_text_box(s, Cm(13), Cm(2.2), Cm(11), Cm(0.6), '全场景业务覆盖', FONT_EN, Pt(18), DARK, True)
scenes = [
    ('顺流交易（母→子）', '未实现利润全部归母'),
    ('逆流交易（子→母）', '少数股东按持股比例分摊'),
    ('平流交易（子↔子）', '按销售方少数股东比例分摊'),
    ('递延所得税', '自动计算（支持税率可调）'),
]
for i, (title, desc) in enumerate(scenes):
    card(s, Cm(13), Cm(3.0) + Cm(2.2) * i, Cm(10.5), Cm(1.9), title, desc)


# ══════════════════════════════════════════════════
# SLIDE 4: 系统架构
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, '系统架构设计 (V2.0 模块化)')

layers = [
    ('表示层  Streamlit UI', '文件上传 | 参数配置 | 可视化仪表盘 | 报表导出', ACCENT),
    ('计算引擎  Engine', '数据校验 → 类型识别 → 利润计算 → 分录生成 → 哈希存证', RGBColor(0x3B, 0x82, 0xF6)),
    ('规则配置  JSON Engine', '三套抵消模板（DOWN/UP/FLAT）· 安全公式求值 · 声明式规则', RGBColor(0x8B, 0x5C, 0xF6)),
    ('数据层  SQLite', '股权配置表 | 交易明细表 | 抵消分录结果表 | 计算日志(calc_log)', GREEN),
]
for i, (title, desc, color) in enumerate(layers):
    y = Cm(2.3) + Cm(2.5) * i
    add_shape(s, Cm(1.5), y, Cm(21.5), Cm(2.0), CARD_BG)
    add_shape(s, Cm(1.5), y, Cm(0.15), Cm(2.0), color)
    add_text_box(s, Cm(2.2), y + Cm(0.2), Cm(20), Cm(0.7), title, FONT_EN, Pt(14), DARK, True)
    add_text_box(s, Cm(2.2), y + Cm(1.0), Cm(20), Cm(0.6), desc, FONT_CN, Pt(11), MUTED)

# 技术栈
add_text_box(s, Cm(1.5), Cm(12.5), Cm(21.5), Cm(0.6),
    '技术栈：Python 3.11+ | Streamlit | pandas | openpyxl | SQLite | Matplotlib | ReportLab | Git',
    FONT_CN, Pt(10.5), BODY, True)


# ══════════════════════════════════════════════════
# SLIDE 5: 核心算法 - 未实现利润计算
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, '核心算法：未实现利润计算机制')

add_text_box(s, Cm(1.2), Cm(2.2), Cm(22), Cm(0.6), '核心公式推导', FONT_EN, Pt(18), DARK, True)

# 公式
add_shape(s, Cm(1.5), Cm(3.0), Cm(21.5), Cm(1.5), CARD_BG)
add_text_box(s, Cm(2.0), Cm(3.2), Cm(20), Cm(1.2),
    '未实现利润 = 销售收入 × 毛利率 × 未对外出售比例\n毛利率 = (销售收入 − 销售成本) ÷ 销售收入    未售比例 = 期末结存数量 ÷ 内部购入总数量',
    FONT_CN, Pt(13), DARK, True)

# 三种场景
add_text_box(s, Cm(1.2), Cm(5.0), Cm(22), Cm(0.6), '三种计算场景', FONT_EN, Pt(16), ACCENT, True)
scene_cards = [
    (Cm(1.5), '全未售', '结存 = 购入 → 全额抵消\n未实现利润 = 收入 − 成本', RED),
    (Cm(8.5), '部分售出', '仅按比例抵消结存未售\n部分的利润', AMBER),
    (Cm(15.5), '全售出', '结存 = 0 → 已实现\n无需抵消，自动跳过', GREEN),
]
for x, title, desc, color in scene_cards:
    add_shape(s, x, Cm(5.8), Cm(6), Cm(2.2), CARD_BG)
    add_shape(s, x, Cm(5.8), Cm(6), Pt(3), color)
    add_text_box(s, x + Cm(0.4), Cm(6.1), Cm(5.2), Cm(0.6), title, FONT_CN, Pt(14), DARK, True)
    add_text_box(s, x + Cm(0.4), Cm(6.8), Cm(5.2), Cm(1.0), desc, FONT_CN, Pt(10.5), MUTED)

# 示例
add_text_box(s, Cm(1.2), Cm(8.5), Cm(22), Cm(0.6), '计算示例（顺流交易）', FONT_EN, Pt(16), ACCENT, True)
add_shape(s, Cm(1.5), Cm(9.2), Cm(21.5), Cm(2.8), DARK_CARD)
example_lines = [
    ('母公司 P 向子公司 S 销售 A 产品 100 件 | 收入 ¥1,000,000 | 成本 ¥600,000 | S 期末对外售出 30 件，剩余 70 件', FONT_CN, Pt(12), WHITE, False),
    ('', FONT_CN, Pt(4), WHITE, False),
    ('① 毛利率 = (100万 − 60万) / 100万 = 40%    ② 未售比 = 70件 / 100件 = 70%', FONT_CN, Pt(13), ACCENT, True),
    ('③ 未实现利润 = 100万 × 40% × 70% = ¥280,000', FONT_CN, Pt(15), WHITE, True),
]
add_multi_text(s, Cm(2.0), Cm(9.4), Cm(20.5), Cm(2.5), example_lines)


# ══════════════════════════════════════════════════
# SLIDE 6: JSON规则引擎
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, '分录自动生成：JSON规则引擎驱动')

add_text_box(s, Cm(1.2), Cm(2.2), Cm(22), Cm(1.0),
    '会计逻辑与程序代码解耦，规则即配置。核心引擎基于三套模板驱动：',
    FONT_CN, Pt(13), BODY)

templates = [
    ('DOWN  顺流交易（母→子）', [
        ('借：营业收入', '贷：营业成本', '贷：存货（未实现利润）'),
        ('借：递延所得税资产', '贷：所得税费用', '不涉及少数股东'),
    ]),
    ('UP  逆流交易（子→母）', [
        ('借：营业收入', '贷：营业成本', '贷：存货（未实现利润）'),
        ('借：递延所得税资产', '贷：所得税费用', ''),
        ('借：少数股东权益', '贷：少数股东损益', '按(1-税率)×持股%分摊'),
    ]),
    ('FLAT  平流交易（子↔子）', [
        ('与逆流同理', '按销售方少数股东', '持股比例分摊'),
    ]),
]
for ti, (title, entries) in enumerate(templates):
    x = Cm(1.5) + Cm(7.5) * ti
    add_shape(s, x, Cm(3.5), Cm(6.8), Cm(5.5), CARD_BG)
    add_shape(s, x, Cm(3.5), Cm(6.8), Pt(3), ACCENT)
    add_text_box(s, x + Cm(0.4), Cm(3.8), Cm(6.0), Cm(0.6), title, FONT_EN, Pt(14), DARK, True)
    for ei, (dr, cr, note) in enumerate(entries):
        y = Cm(4.6) + Cm(1.3) * ei
        add_text_box(s, x + Cm(0.4), y, Cm(6.0), Cm(0.5), f'{dr}    {cr}', FONT_CN, Pt(11), BODY)
        if note:
            add_text_box(s, x + Cm(0.4), y + Cm(0.5), Cm(6.0), Cm(0.5), note, FONT_CN, Pt(9.5), MUTED)

add_text_box(s, Cm(1.5), Cm(9.5), Cm(21.5), Cm(2),
    '设计理念：新增抵消场景（如固定资产内部交易）只需在 rules.json 中增加配置条目，核心引擎代码无需修改。'
    '\n所有分录基于安全命名空间求值，变量受限（revenue / cost / unrealized_profit / tax_rate / minority_ratio），防止代码注入。',
    FONT_CN, Pt(12), MUTED)


# ══════════════════════════════════════════════════
# SLIDE 7: 数据可视化与追溯
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, '数据可视化与穿透式追溯')

viz = [
    ('动态分析仪表盘', '5项KPI（交易额/利润/抵消笔数/存货/净利影响）+ 抵消金额构成柱状图\n+ 交易类型分布图 + 抵消前后对比表。支持按交易类型标签过滤分录。'),
    ('自动试算平衡', '自动汇总借方合计与贷方合计，差额为 0 显示"✅ 借贷平衡"。\n支持人工干预：修改金额、删除分录行、追加入工分录（标注"✏️ 人工追加"）。'),
    ('SHA-256 哈希链追溯', '每批次计算生成完整 JSON 日志。每笔分录哈希含上一笔哈希值，\n形成防篡改哈希链。预留 T33 区块链存证接口，赋能穿透式监管。'),
]
for i, (title, desc) in enumerate(viz):
    y = Cm(2.3) + Cm(3.5) * i
    add_shape(s, Cm(1.5), y, Cm(21.5), Cm(3.0), CARD_BG)
    add_shape(s, Cm(1.5), y, Pt(3), Cm(3.0), ACCENT)
    add_text_box(s, Cm(2.0), y + Cm(0.3), Cm(20), Cm(0.6), title, FONT_EN, Pt(15), DARK, True)
    add_text_box(s, Cm(2.0), y + Cm(1.1), Cm(20), Cm(1.8), desc, FONT_CN, Pt(11), MUTED)


# ══════════════════════════════════════════════════
# SLIDE 8: 与老师平台对比
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, '系统对标：与现有DRP T11的对比')

add_text_box(s, Cm(1.2), Cm(2.2), Cm(22), Cm(0.8),
    '老师"集团合并报表平台"（19.group-consolidation）的"合并抵消处理"页面与本系统功能对比：',
    FONT_CN, Pt(12), BODY)

# 对比表
table_data = [
    ['对比维度', '现有 DRP T11 页面', '本系统 V2.0', '优势'],
    ['抵消分录', '空表单，需手动填写', 'AI规则引擎自动生成', '✅ 零人工'],
    ['未实现利润计算', '无自动计算能力', '毛利率×未售比例自动计算', '✅ 自动化'],
    ['递延所得税', '无', '自动计算（税率可调）', '✅ 智能'],
    ['少数股东分摊', '无', '逆流/平流自动分摊', '✅ 完整'],
    ['试算平衡', '有独立表单', '内嵌自动验证，差额即时展示', '✅ 即时'],
    ['历史追溯', '基础处理历史记录', '批次分组 + 哈希链 + JSON日志', '✅ 防篡改'],
    ['分录调整', '表单增删改', '金额修改/行删除/人工追加', '➡ 持平'],
]
for ri, row_data in enumerate(table_data):
    y = Cm(3.2) + Cm(1.1) * ri
    is_header = ri == 0
    for ci, cell_text in enumerate(row_data):
        widths = [Cm(3.5), Cm(7.2), Cm(7.5), Cm(4)]
        x = Cm(1.5) + sum(Cm(0) if i == 0 else widths[i-1] for i in range(ci + 1))
        actual_x = Cm(1.5)
        if ci == 1: actual_x = Cm(5.0)
        elif ci == 2: actual_x = Cm(12.2)
        elif ci == 3: actual_x = Cm(19.7)
        
        if is_header:
            add_shape(s, actual_x, y, widths[ci], Cm(0.9), ACCENT)
            add_text_box(s, actual_x + Cm(0.15), y + Cm(0.1), widths[ci] - Cm(0.3), Cm(0.7), cell_text, FONT_CN, Pt(10), WHITE, True)
        else:
            bg = CARD_BG if ri % 2 == 0 else BG
            add_shape(s, actual_x, y, widths[ci], Cm(0.9), bg)
            color = GREEN if '✅' in cell_text else (BODY if '➡' in cell_text else BODY)
            add_text_box(s, actual_x + Cm(0.15), y + Cm(0.1), widths[ci] - Cm(0.3), Cm(0.7), cell_text, FONT_CN, Pt(9.5), color)


# ══════════════════════════════════════════════════
# SLIDE 9: 交付范围与未来
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'V2.0 交付范围与未来演进')

# 本期交付
add_text_box(s, Cm(1.2), Cm(2.2), Cm(11), Cm(0.6), 'V2.0 交付物清单（本期）', FONT_EN, Pt(18), DARK, True)
deliverables = [
    '✅ 系统源代码：Python 3.11+ 模块化架构',
    '✅ 三表上传 + 自动计算 + 分录生成 + 试算平衡',
    '✅ 可视化仪表盘 + 图表 + 抵消前后对比',
    '✅ 分录筛选 + 手动调整 + 历史追溯 + 哈希链',
    '✅ Excel/PDF 双格式导出',
    '📘 系统需求说明书 V4.0（GB/T 9385）',
    '📗 用户操作手册 V2.0',
    '📊 介绍PPT V3.0',
    '📋 Excel模板 × 3 + 预设演示数据',
]
for i, item in enumerate(deliverables):
    add_text_box(s, Cm(1.5), Cm(3.0) + Cm(0.8) * i, Cm(11), Cm(0.7), item, FONT_CN, Pt(11), BODY)

# 未来扩展
add_text_box(s, Cm(13), Cm(2.2), Cm(11), Cm(0.6), 'V3.0+ 未来扩展', FONT_EN, Pt(18), DARK, True)
future = [
    '业务：加权平均法 / 跌价准备 / 连续编报滚调',
    '架构：Java Spring Boot + Vue 3 重构',
    '信创：国产数据库（达梦/人大金仓）适配',
    '安全：对接 T33 区块链存证上链',
    '集成：RESTful API 对接 T10/T05/T21',
    '合规：等保三级 + 多会计准则切换',
]
for i, item in enumerate(future):
    add_text_box(s, Cm(13.5), Cm(3.0) + Cm(0.8) * i, Cm(10), Cm(0.7), f'○  {item}', FONT_CN, Pt(11), ACCENT, True)


# ══════════════════════════════════════════════════
# SLIDE 10: 总结
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_shape(s, Cm(1.5), Cm(3.0), Cm(3), Pt(3), ACCENT)

add_multi_text(s, Cm(1.5), Cm(3.5), Cm(22), Cm(3), [
    ('感谢聆听', FONT_EN, Pt(39), WHITE, True),
])
add_text_box(s, Cm(1.5), Cm(6.0), Cm(22), Cm(1.5),
    '基于AI规则引擎的内部存货交易未实现利润自动抵消系统\nDRP平台 T11 合并报表与一键出表 — 子模块',
    FONT_CN, Pt(14), SUBTLE)

add_multi_text(s, Cm(1.5), Cm(10.0), Cm(22), Cm(3), [
    ('中国社会科学院大学 · 2025级MPAcc二班', FONT_CN, Pt(15), WHITE, True),
    ('小组成员：尚浩然、李金阳 | 指导老师：张金昌', FONT_CN, Pt(13), SUBTLE, False),
    ('在线体验：t11-consolidation-ucassmpacc2025.streamlit.app', FONT_CN, Pt(11), MUTED, False),
])


# ── 保存 ──
output = Path(__file__).parent / 'T11-介绍PPT-v3.pptx'
prs.save(str(output))
print(f'✅ PPT已生成：{output}')
print(f'   共 {len(prs.slides)} 页')
